"""
KRISH PDF Toolkit - Complete Processing Engine
"""

import fitz
import os
import io
import zipfile
import tempfile
import subprocess
import platform
import json
from pathlib import Path
from PIL import Image

from pdf2docx import Converter as Pdf2DocxConverter
import pdfplumber
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

try:
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False


def hex_to_rgb(hex_color):
    if not hex_color or hex_color in ("", "none", "transparent"):
        return None
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return (0, 0, 0)
    r = int(hex_color[0:2], 16) / 255.0
    g = int(hex_color[2:4], 16) / 255.0
    b = int(hex_color[4:6], 16) / 255.0
    return (r, g, b)


FONT_MAP = {
    "Helvetica": "helv",
    "Helvetica-Bold": "hebo",
    "Times-Roman": "tiro",
    "Times-Bold": "tibo",
    "Courier": "cour",
    "Courier-Bold": "cobo",
}


class PDFEngine:

    @staticmethod
    def merge(pdf_paths, output_path):
        result = fitz.open()
        for path in pdf_paths:
            doc = fitz.open(path)
            result.insert_pdf(doc)
            doc.close()
        result.save(output_path)
        result.close()

    @staticmethod
    def split(pdf_path, mode="all", pages=None, every_n=1):
        doc = fitz.open(pdf_path)
        output_dir = tempfile.mkdtemp()
        output_files = []
        if mode == "all":
            for i in range(len(doc)):
                out = fitz.open()
                out.insert_pdf(doc, from_page=i, to_page=i)
                p = os.path.join(output_dir, f"page_{i+1}.pdf")
                out.save(p); out.close(); output_files.append(p)
        elif mode == "ranges" and pages:
            for idx, (start, end) in enumerate(pages):
                out = fitz.open()
                out.insert_pdf(doc, from_page=start-1, to_page=end-1)
                p = os.path.join(output_dir, f"pages_{start}-{end}.pdf")
                out.save(p); out.close(); output_files.append(p)
        elif mode == "every_n":
            total = len(doc)
            for i in range(0, total, every_n):
                out = fitz.open()
                end = min(i + every_n - 1, total - 1)
                out.insert_pdf(doc, from_page=i, to_page=end)
                p = os.path.join(output_dir, f"pages_{i+1}-{end+1}.pdf")
                out.save(p); out.close(); output_files.append(p)
        elif mode == "extract" and pages:
            out = fitz.open()
            for pg in pages:
                out.insert_pdf(doc, from_page=pg-1, to_page=pg-1)
            p = os.path.join(output_dir, "extracted_pages.pdf")
            out.save(p); out.close(); output_files.append(p)
        doc.close()
        if len(output_files) > 1:
            zip_path = os.path.join(output_dir, "split_pdfs.zip")
            with zipfile.ZipFile(zip_path, "w") as zf:
                for f in output_files:
                    zf.write(f, os.path.basename(f))
            return zip_path
        return output_files[0] if output_files else pdf_path

    @staticmethod
    def reorder(pdf_path, new_order, output_path):
        doc = fitz.open(pdf_path)
        doc.select([p - 1 for p in new_order])
        doc.save(output_path); doc.close()

    @staticmethod
    def pdf_to_word(pdf_path, output_path):
        cv = Pdf2DocxConverter(pdf_path)
        cv.convert(output_path); cv.close()

    @staticmethod
    def pdf_to_excel(pdf_path, output_path):
        wb = Workbook(); ws = wb.active; ws.title = "PDF Tables"
        row_offset = 1
        with pdfplumber.open(pdf_path) as pdf:
            for pn, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        for row in table:
                            for ci, cell in enumerate(row):
                                ws.cell(row=row_offset, column=ci+1,
                                        value=cell or "")
                            row_offset += 1
                        row_offset += 1
                else:
                    text = page.extract_text()
                    if text:
                        ws.cell(row=row_offset, column=1,
                                value=f"--- Page {pn+1} ---")
                        row_offset += 1
                        for line in text.split("\n"):
                            ws.cell(row=row_offset, column=1, value=line)
                            row_offset += 1
                        row_offset += 1
        wb.save(output_path)

    @staticmethod
    def pdf_to_ppt(pdf_path, output_path, dpi=200):
        doc = fitz.open(pdf_path)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        for pn in range(len(doc)):
            page = doc[pn]
            pix = page.get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72))
            img_stream = io.BytesIO(pix.tobytes("png"))
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0),
                                     prs.slide_width, prs.slide_height)
        doc.close(); prs.save(output_path)

    @staticmethod
    def office_to_pdf(input_path, output_dir):
        system = platform.system()
        if system == "Windows":
            candidates = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            ]
            lo_cmd = next((p for p in candidates if os.path.exists(p)), None)
            if lo_cmd is None:
                raise FileNotFoundError("LibreOffice not found!")
        elif system == "Darwin":
            lo_cmd = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        else:
            lo_cmd = "libreoffice"
        cmd = [lo_cmd, "--headless", "--convert-to", "pdf",
               "--outdir", output_dir, input_path]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice error: {result.stderr}")
        base = os.path.splitext(os.path.basename(input_path))[0]
        return os.path.join(output_dir, f"{base}.pdf")

    @staticmethod
    def add_watermark(pdf_path, output_path, text="CONFIDENTIAL",
                      fontsize=60, opacity=0.3, rotation=45,
                      font_name="Helvetica"):
        doc = fitz.open(pdf_path)
        c = max(0.0, min(1.0, 1.0 - opacity))
        color = (c, c, c)
        fitz_font = FONT_MAP.get(font_name, "helv")
        for page in doc:
            rect = page.rect
            cx, cy = rect.width / 2, rect.height / 2
            text_len = fitz.get_text_length(text, fontname=fitz_font,
                                            fontsize=fontsize)
            point = fitz.Point(cx - text_len / 2, cy)
            morph = (fitz.Point(cx, cy), fitz.Matrix(rotation))
            page.insert_text(point, text, fontname=fitz_font,
                             fontsize=fontsize, color=color,
                             morph=morph, overlay=True)
        doc.save(output_path); doc.close()

    @staticmethod
    def add_page_numbers(pdf_path, output_path, position="bottom-center",
                         start_num=1, fontsize=12, format_str="{n}"):
        doc = fitz.open(pdf_path)
        total = len(doc); margin = 36
        for i, page in enumerate(doc):
            rect = page.rect; num = start_num + i
            text = format_str.replace("{n}", str(num)).replace("{total}", str(total))
            tl = fitz.get_text_length(text, fontsize=fontsize)
            positions = {
                "bottom-center": fitz.Point(rect.width/2-tl/2, rect.height-margin),
                "bottom-left": fitz.Point(margin, rect.height-margin),
                "bottom-right": fitz.Point(rect.width-margin-tl, rect.height-margin),
                "top-center": fitz.Point(rect.width/2-tl/2, margin+fontsize),
                "top-left": fitz.Point(margin, margin+fontsize),
                "top-right": fitz.Point(rect.width-margin-tl, margin+fontsize),
            }
            pt = positions.get(position, positions["bottom-center"])
            page.insert_text(pt, text, fontsize=fontsize, color=(0, 0, 0))
        doc.save(output_path); doc.close()

    @staticmethod
    def add_header_footer(pdf_path, output_path, header="", footer="",
                          fontsize=10):
        doc = fitz.open(pdf_path); margin = 36
        for page in doc:
            rect = page.rect
            if header:
                hl = fitz.get_text_length(header, fontsize=fontsize)
                page.insert_text(fitz.Point(rect.width/2-hl/2, margin),
                                 header, fontsize=fontsize, color=(0, 0, 0))
            if footer:
                fl = fitz.get_text_length(footer, fontsize=fontsize)
                page.insert_text(fitz.Point(rect.width/2-fl/2,
                                            rect.height-margin/2),
                                 footer, fontsize=fontsize, color=(0, 0, 0))
        doc.save(output_path); doc.close()

    @staticmethod
    def ocr_pdf(pdf_path, output_path, language="eng"):
        if not OCR_AVAILABLE:
            raise ImportError("pytesseract not installed")
        doc = fitz.open(pdf_path); result = fitz.open()
        for pn in range(len(doc)):
            page = doc[pn]
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            ocr_bytes = pytesseract.image_to_pdf_or_hocr(
                img, lang=language, extension="pdf")
            ocr_doc = fitz.open("pdf", ocr_bytes)
            result.insert_pdf(ocr_doc); ocr_doc.close()
        doc.close(); result.save(output_path); result.close()

    @staticmethod
    def redact_text(pdf_path, output_path, search_texts):
        doc = fitz.open(pdf_path)
        for page in doc:
            for text in search_texts:
                hits = page.search_for(text)
                for rect in hits:
                    page.add_redact_annot(rect, fill=(0, 0, 0))
            page.apply_redactions()
        doc.save(output_path); doc.close()

    @staticmethod
    def repair_pdf(pdf_path, output_path):
        doc = fitz.open(pdf_path)
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        doc.close()

    @staticmethod
    def compress(pdf_path, output_path, quality="medium"):
        doc = fitz.open(pdf_path)
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        doc.close()

    @staticmethod
    def rotate(pdf_path, output_path, angle=90, page_numbers=None):
        doc = fitz.open(pdf_path); total = len(doc)
        targets = range(total) if page_numbers is None else \
            [p-1 for p in page_numbers if 0 < p <= total]
        for i in targets:
            doc[i].set_rotation((doc[i].rotation + angle) % 360)
        doc.save(output_path); doc.close()

    @staticmethod
    def pdf_to_images(pdf_path, img_format="png", dpi=200):
        doc = fitz.open(pdf_path)
        output_dir = tempfile.mkdtemp(); output_files = []
        for i in range(len(doc)):
            pix = doc[i].get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72))
            fp = os.path.join(output_dir, f"page_{i+1}.{img_format}")
            if img_format in ("jpg", "jpeg"):
                pix.save(fp, jpg_quality=95)
            else:
                pix.save(fp)
            output_files.append(fp)
        doc.close()
        zip_path = os.path.join(output_dir, "pdf_images.zip")
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for f in output_files:
                zf.write(f, os.path.basename(f))
        return zip_path

    @staticmethod
    def images_to_pdf(image_paths, output_path):
        doc = fitz.open()
        for img_path in image_paths:
            img = fitz.open(img_path)
            pdf_bytes = img.convert_to_pdf(); img.close()
            img_pdf = fitz.open("pdf", pdf_bytes)
            doc.insert_pdf(img_pdf); img_pdf.close()
        doc.save(output_path); doc.close()

    @staticmethod
    def protect(pdf_path, output_path, user_password, owner_password=None):
        if owner_password is None:
            owner_password = user_password
        doc = fitz.open(pdf_path)
        perm = int(fitz.PDF_PERM_PRINT | fitz.PDF_PERM_COPY
                    | fitz.PDF_PERM_ANNOTATE)
        doc.save(output_path, encryption=fitz.PDF_ENCRYPT_AES_256,
                 owner_pw=owner_password, user_pw=user_password,
                 permissions=perm)
        doc.close()

    @staticmethod
    def unlock(pdf_path, output_path, password):
        doc = fitz.open(pdf_path)
        if doc.is_encrypted:
            if not doc.authenticate(password):
                doc.close()
                raise ValueError("Wrong password!")
        doc.save(output_path, encryption=fitz.PDF_ENCRYPT_NONE)
        doc.close()

    @staticmethod
    def delete_pages(pdf_path, output_path, pages_to_delete):
        """Delete specific pages. pages_to_delete is 1-based list."""
        doc = fitz.open(pdf_path)
        total = len(doc)
        valid = sorted(
            [p - 1 for p in pages_to_delete if 0 < p <= total],
            reverse=True
        )
        if len(valid) >= total:
            doc.close()
            raise ValueError("Cannot delete ALL pages!")
        for page_idx in valid:
            doc.delete_page(page_idx)
        doc.save(output_path)
        doc.close()

    @staticmethod
    def apply_edits(pdf_path, output_path, annotations, dpi=150):
        doc = fitz.open(pdf_path)
        scale = 72.0 / dpi
        for ann in annotations:
            page_num = ann.get("page", 0)
            if page_num >= len(doc):
                continue
            page = doc[page_num]
            x = ann.get("x", 0) * scale
            y = ann.get("y", 0) * scale
            ann_type = ann.get("type", "")
            color = hex_to_rgb(ann.get("color", "#000000"))
            if ann_type == "text":
                fs = ann.get("fontSize", 14) * scale
                font_key = FONT_MAP.get(
                    ann.get("fontFamily", "Helvetica"), "helv")
                page.insert_text(
                    fitz.Point(x, y + fs), ann.get("text", ""),
                    fontname=font_key, fontsize=fs,
                    color=color or (0, 0, 0))
            elif ann_type == "rect":
                w = ann.get("width", 100) * scale
                h = ann.get("height", 50) * scale
                rect = fitz.Rect(x, y, x + w, y + h)
                shape = page.new_shape(); shape.draw_rect(rect)
                fill = hex_to_rgb(ann.get("fill", ""))
                lw = ann.get("lineWidth", 2) * scale
                shape.finish(color=color, fill=fill, width=lw)
                shape.commit()
            elif ann_type == "circle":
                w = ann.get("width", 100) * scale
                h = ann.get("height", 100) * scale
                rect = fitz.Rect(x, y, x + w, y + h)
                shape = page.new_shape(); shape.draw_oval(rect)
                fill = hex_to_rgb(ann.get("fill", ""))
                lw = ann.get("lineWidth", 2) * scale
                shape.finish(color=color, fill=fill, width=lw)
                shape.commit()
            elif ann_type == "freehand":
                points = ann.get("points", [])
                if len(points) < 2:
                    continue
                shape = page.new_shape()
                lw = ann.get("lineWidth", 2) * scale
                for i in range(1, len(points)):
                    p1 = fitz.Point(points[i-1][0]*scale, points[i-1][1]*scale)
                    p2 = fitz.Point(points[i][0]*scale, points[i][1]*scale)
                    shape.draw_line(p1, p2)
                shape.finish(color=color, width=lw); shape.commit()
            elif ann_type == "image":
                w = ann.get("width", 100) * scale
                h = ann.get("height", 100) * scale
                img_data = ann.get("imageData", "")
                if img_data:
                    import base64
                    if "," in img_data:
                        img_data = img_data.split(",")[1]
                    img_bytes = base64.b64decode(img_data)
                    rect = fitz.Rect(x, y, x + w, y + h)
                    page.insert_image(rect, stream=img_bytes)
        doc.save(output_path); doc.close()

    @staticmethod
    def bates_number(pdf_path, output_path, prefix="DOC",
                     start=1, digits=6, position="bottom-right"):
        doc = fitz.open(pdf_path); margin = 36; fontsize = 10
        for i, page in enumerate(doc):
            rect = page.rect
            num_str = f"{prefix}{str(start+i).zfill(digits)}"
            tl = fitz.get_text_length(num_str, fontsize=fontsize)
            positions = {
                "bottom-right": fitz.Point(rect.width-margin-tl,
                                           rect.height-margin),
                "bottom-left": fitz.Point(margin, rect.height-margin),
                "top-right": fitz.Point(rect.width-margin-tl, margin+fontsize),
                "top-left": fitz.Point(margin, margin+fontsize),
            }
            pt = positions.get(position, positions["bottom-right"])
            page.insert_text(pt, num_str, fontsize=fontsize, color=(0, 0, 0))
        doc.save(output_path); doc.close()

    @staticmethod
    def get_preview(pdf_path, page_num=0, dpi=150):
        doc = fitz.open(pdf_path)
        page_num = min(page_num, len(doc) - 1)
        pix = doc[page_num].get_pixmap(
            matrix=fitz.Matrix(dpi/72, dpi/72))
        img_bytes = pix.tobytes("png")
        doc.close()
        return img_bytes

    @staticmethod
    def get_page_count(pdf_path):
        doc = fitz.open(pdf_path); n = len(doc); doc.close()
        return n

    @staticmethod
    def get_metadata(pdf_path):
        doc = fitz.open(pdf_path)
        meta = doc.metadata or {}
        info = {"pages": len(doc),
                "size_bytes": os.path.getsize(pdf_path),
                "size_mb": round(os.path.getsize(pdf_path)/(1024*1024), 2)}
        info.update(meta); doc.close()
        return info